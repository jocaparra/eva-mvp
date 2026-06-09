"""Gera PPT a partir de template .pptx do cliente (python-pptx)."""

from __future__ import annotations

import os
from datetime import date
from typing import Any, Dict, Iterable, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _join_items(items: Any) -> str:
    if not items:
        return ""
    if isinstance(items, list):
        return ", ".join(str(x) for x in items)
    return str(items)


def build_shape_values(state: dict) -> Dict[str, str]:
    """Mapeia shape.name → texto para substituição no template."""
    r = state.get("research_structured") or {}
    f = state.get("financial_structured") or {}
    company = r.get("company_name") or state.get("company_name", "Empresa")

    comparables = f.get("comparable_companies") or []
    comp_text = "\n".join(
        f"{c.get('name', 'N/A')}: EV/EBITDA {c.get('ev_ebitda', 'est.')}, P/L {c.get('pe', 'est.')}"
        for c in comparables[:5]
    )

    red_flags = f.get("red_flags") or []
    if isinstance(red_flags, list):
        risks = "\n".join(f"• {flag}" for flag in red_flags[:5])
    else:
        risks = str(red_flags)

    risks_list = f.get("risks") or []
    if isinstance(risks_list, list) and risks_list:
        risk_lines = []
        for item in risks_list[:5]:
            if isinstance(item, dict):
                risk_lines.append(
                    f"• {item.get('risco', '')} ({item.get('nivel', '')})"
                )
            else:
                risk_lines.append(f"• {item}")
        if risk_lines:
            risks = "\n".join(risk_lines)

    return {
        "company_name": str(company),
        "tagline": str(r.get("tagline") or ""),
        "description": str(r.get("description") or ""),
        "sector": str(r.get("sector") or ""),
        "founded": str(r.get("founded") or ""),
        "headquarters": str(r.get("headquarters") or ""),
        "founders": str(r.get("founders") or ""),
        "employees": str(r.get("employees") or ""),
        "business_model": str(r.get("business_model") or ""),
        "market_position": str(r.get("market_position") or ""),
        "recent_news": str(r.get("recent_news") or ""),
        "key_products": _join_items(r.get("key_products")),
        "main_competitors": _join_items(r.get("main_competitors")),
        "revenue_2022": str(f.get("revenue_2022") or ""),
        "revenue_2023": str(f.get("revenue_2023") or ""),
        "revenue_2024": str(f.get("revenue_2024") or ""),
        "ebitda": str(f.get("ebitda") or ""),
        "net_margin": str(f.get("net_margin") or ""),
        "growth_yoy": str(f.get("growth_yoy") or ""),
        "valuation": str(f.get("valuation") or ""),
        "valuation_low": str(f.get("valuation_low") or ""),
        "valuation_mid": str(f.get("valuation_mid") or ""),
        "valuation_high": str(f.get("valuation_high") or ""),
        "ev_ebitda": str(f.get("ev_ebitda") or ""),
        "pe_ratio": str(f.get("pe_ratio") or ""),
        "comparables": comp_text or "Comparáveis do setor est.",
        "risks": risks or str(f.get("risk_score") or "Riscos em análise est."),
        "risk_score": str(f.get("risk_score") or ""),
    }


def _iter_shapes(shapes) -> Iterable:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _set_shape_text(shape, text: str) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    if not tf.paragraphs:
        return False
    tf.paragraphs[0].text = text
    for para in tf.paragraphs[1:]:
        para.text = ""
    return True


def fill_template(template_path: str, values: Dict[str, str]) -> Presentation:
    prs = Presentation(template_path)
    replaced: set[str] = set()

    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            name = (shape.name or "").strip()
            if not name or name not in values:
                continue
            text = values[name]
            if text and _set_shape_text(shape, text):
                replaced.add(name)

    print(f"document_pptx: shapes preenchidos: {sorted(replaced)}")
    return prs


def append_citations_slide(prs: Presentation, state: dict) -> None:
    """Adiciona slide de referências com citações do data room e fontes externas."""
    citations = (
        (state.get("financial_citations") or [])
        + (state.get("research_citations") or [])
    )
    if not citations:
        return

    seen: set[str] = set()
    lines: list[str] = ["Fontes e referências:"]
    for idx, citation in enumerate(citations, start=1):
        key = citation.get("chunk_id") or f"{citation.get('source_file')}-{idx}"
        if key in seen:
            continue
        seen.add(key)
        source = citation.get("source") or "data_room"
        src_file = citation.get("source_file") or "—"
        page = citation.get("page") or "—"
        quote = (citation.get("quote") or "")[:120]
        lines.append(f"[{idx}] ({source}) {src_file}, pág. {page}: {quote}")

    slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Referências"
    body = "\n".join(lines[:12])
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            shape.text_frame.text = body
            break


def generate_pptx_from_template(
    state: dict,
    template_path: str,
    output_path: str,
) -> dict:
    doc_type = state.get("document_type", "CIM")
    company = state.get("research_structured", {}).get("company_name") or state.get(
        "company_name", "empresa"
    )
    job_id = state.get("job_id", "unknown")

    values = build_shape_values(state)
    prs = fill_template(template_path, values)
    append_citations_slide(prs, state)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    filename = f"{company}_{doc_type}_{date.today()}.pptx"
    return {
        "output_path": output_path,
        "output_filename": filename,
        "ppt_path": output_path,
        "ppt_filename": filename,
        "job_id": job_id,
    }


def generate_from_client_template(state: dict) -> dict:
    """Gera PPT usando template do client_id (telefone) via Supabase Storage."""
    from app.utils.template import cleanup_temp_template, get_template_path

    client_id = state.get("client_id", "default")
    job_id = state.get("job_id", "unknown")
    template_path = get_template_path(client_id, job_id=job_id)
    try:
        if not template_path:
            raise FileNotFoundError(f"Template não encontrado para {client_id}")

        output_path = f"outputs/{job_id}.pptx"
        return generate_pptx_from_template(state, template_path, output_path)
    finally:
        if template_path:
            cleanup_temp_template(template_path)
