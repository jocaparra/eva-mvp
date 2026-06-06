"""Extrai texto de PDF/DOCX/PPTX/XLSX preservando número de página/planilha."""

from __future__ import annotations

import io
from dataclasses import dataclass
from typing import List, Optional


@dataclass(frozen=True)
class TextSegment:
    """Trecho de texto com referência de página ou slide/planilha."""

    page: int
    text: str


def _detect_type(filename: str, mime_type: str = "") -> Optional[str]:
    name = (filename or "").lower()
    mime = (mime_type or "").lower()
    if name.endswith(".pdf") or mime == "application/pdf":
        return "pdf"
    if name.endswith(".docx") or "wordprocessingml" in mime:
        return "docx"
    if name.endswith(".pptx") or "presentationml" in mime:
        return "pptx"
    if name.endswith((".xlsx", ".xls")) or "spreadsheetml" in mime or mime.endswith("/ms-excel"):
        return "excel"
    return None


def load_pdf(data: bytes) -> List[TextSegment]:
    import fitz

    doc = fitz.open(stream=data, filetype="pdf")
    segments: List[TextSegment] = []
    try:
        for idx, page in enumerate(doc, start=1):
            text = page.get_text().strip()
            if text:
                segments.append(TextSegment(page=idx, text=text))
    finally:
        doc.close()
    return segments


def load_docx(data: bytes) -> List[TextSegment]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    body = "\n".join(parts).strip()
    return [TextSegment(page=1, text=body)] if body else []


def load_pptx(data: bytes) -> List[TextSegment]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    segments: List[TextSegment] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                texts.append(text)
        if texts:
            segments.append(TextSegment(page=idx, text="\n".join(texts)))
    return segments


def load_excel(data: bytes, filename: str = "") -> List[TextSegment]:
    name = (filename or "").lower()
    if name.endswith(".xls") and not name.endswith(".xlsx"):
        return _load_xls(data)
    return _load_xlsx(data)


def _load_xlsx(data: bytes) -> List[TextSegment]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    segments: List[TextSegment] = []
    try:
        for sheet_idx, sheet in enumerate(wb.worksheets, start=1):
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                header = f"Planilha: {sheet.title}"
                segments.append(TextSegment(page=sheet_idx, text=header + "\n" + "\n".join(rows)))
    finally:
        wb.close()
    return segments


def _load_xls(data: bytes) -> List[TextSegment]:
    import xlrd

    book = xlrd.open_workbook(file_contents=data)
    segments: List[TextSegment] = []
    for sheet_idx in range(book.nsheets):
        sheet = book.sheet_by_index(sheet_idx)
        rows: list[str] = []
        for row_idx in range(sheet.nrows):
            cells = [str(sheet.cell_value(row_idx, col_idx)).strip() for col_idx in range(sheet.ncols)]
            cells = [c for c in cells if c]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            header = f"Planilha: {sheet.name}"
            segments.append(TextSegment(page=sheet_idx + 1, text=header + "\n" + "\n".join(rows)))
    return segments


def load_document(data: bytes, filename: str, mime_type: str = "") -> List[TextSegment]:
    """Carrega documento e retorna segmentos com metadado de página/planilha/slide."""
    doc_type = _detect_type(filename, mime_type)
    if doc_type == "pdf":
        segments = load_pdf(data)
    elif doc_type == "docx":
        segments = load_docx(data)
    elif doc_type == "pptx":
        segments = load_pptx(data)
    elif doc_type == "excel":
        segments = load_excel(data, filename)
    else:
        raise ValueError(f"Tipo não suportado para ingestão: {filename or mime_type}")

    if not segments:
        raise ValueError("Não foi possível extrair texto do documento.")
    return segments


def is_ingestible(filename: str, mime_type: str = "") -> bool:
    """Indica se o arquivo pode ser ingerido no data room."""
    return _detect_type(filename, mime_type) is not None
