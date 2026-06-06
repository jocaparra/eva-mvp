"""Lógica PPTX compartilhada (template, Presenton, slide de referências)."""

from __future__ import annotations

import os
from typing import Optional, Tuple

from app.generation.citations import CitationBundle, collect_citations


def append_citations_from_bundle(prs, bundle: CitationBundle) -> None:
    """Slide de referências usando camada compartilhada de citações."""
    if not bundle.has_citations:
        return

    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = bundle.references_heading()
    body = bundle.references_body()[:3000]
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape != slide.shapes.title:
            shape.text_frame.text = body
            break


def render_pptx(deal_state: dict, output_path: str) -> Tuple[str, Optional[str]]:
    """Gera PPTX via template cliente, fallback Presenton."""
    from app.agents.document import presenton_document

    client_id = deal_state.get("client_id", "default")
    job_id = deal_state.get("job_id", "unknown")
    company = deal_state.get("company_name", "empresa")
    doc_type = deal_state.get("document_type", "CIM")

    bundle = collect_citations(deal_state)
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    from app.utils.template import cleanup_temp_template, get_template_path

    template_path = get_template_path(client_id, job_id=job_id)
    if template_path and os.path.exists(template_path):
        try:
            from app.agents.document_pptx import fill_template, build_shape_values

            values = build_shape_values(deal_state)
            from pptx import Presentation

            prs = fill_template(template_path, values)
            append_citations_from_bundle(prs, bundle)
            prs.save(output_path)
            return output_path, None
        except Exception as exc:
            print(f"[cim_pptx] template falhou, fallback Presenton: {exc}")
        finally:
            cleanup_temp_template(template_path)

    # Presenton usa conteúdo plain pré-montado quando disponível
    if deal_state.get("_generation_plain"):
        deal_state = dict(deal_state)
        deal_state["document_type"] = doc_type

    result = presenton_document(deal_state)
    path = result.get("ppt_path") or output_path
    # Presenton salva em outputs/{job_id}.pptx — garantir citações se template local falhou
    if path and os.path.isfile(path) and bundle.has_citations:
        try:
            from pptx import Presentation

            prs = Presentation(path)
            append_citations_from_bundle(prs, bundle)
            prs.save(path)
        except Exception as exc:
            print(f"[cim_pptx] não foi possível anexar referências: {exc}")
    return path, result.get("edit_url")
